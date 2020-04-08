#测试下载 百度文库的ppt
#
#--------------------------------篮子专属    -----------------
#
#
#
#测试链接：https://wenku.baidu.com/view/1c691781ff00bed5b8f31d48?from=search
#
#
#
#测试：第一次，打算用自动化模块获取ppt里面的图片地址   ok
#测试：第二次，尝试把下载的图片加入ppt模板中           ok
#测试：第三次，尝试加入  爬取 word 文档      好像只有一半（不完整）
#测试：第4次， 加入点击继续阅读    ok
#测试：第5次，尝试加入  爬取 word 文档    格式问题
#
#
#
#wordc测试链接https://wenku.baidu.com/view/b0db8c8d6d85ec3a87c24028915f804d2a16873a.html?from=search
##================================d分割线=========================================================
#模块引入：
import os
import re
import sys
import time
import requests
from lxml import etree
from selenium import webdriver

from pptx import Presentation  # 生成ppt需要的包
from pptx.util import Inches

#获取网页链接源代码的函数
def get_request(url):
        try:
                
                ##创建chrome参数对象
                opt = webdriver.ChromeOptions()
                # # 把chrome设置成无界面模式，不论windows还是linux都可以，自动适配对应参数
                opt.set_headless()
                # # 创建chrome无界面对象
                driver = webdriver.Chrome(options=opt)	
                ##打开chrome浏览器
                #driver = webdriver.Chrome()#有界面。可视
                #打开浏览器，输入网址
                driver.get(url)
                time.sleep(5)
                #这里曾经尝试点击  还有多少页继续阅读   但是定位失败，点击不了
                #driver.find_element_by_class_name('banner-more-btn').click()
                #driver.find_element_by_xpath('//*[@id="html-reader-go-more"]/div[2]/div[1]/p').click()
                #这里曾经尝试移动鼠标位置，ppt 的ok    word 的元素定位失败
                 # 点击继续阅读,来源网络不东
                hidden_div = driver.find_element_by_xpath("//div[@id='html-reader-go-more']")
                got_btn = driver.find_element_by_xpath("//div[@class='banner-more-btn']/span")
                xy=driver.find_element_by_xpath('//*[@id="ft"]/div')

                actions = webdriver.ActionChains(driver)
                actions.move_to_element(hidden_div)
                time.sleep(1)

                actions.click(got_btn)
                
                time.sleep(1)
                actions.move_to_element(xy)
                time.sleep(2)

                actions.perform()
        except:
                 return driver.page_source
                
        
        return driver.page_source



#ppt  下载图片，文件标题
def down_jpg(request):
	#利用正则表达式匹配 
	img_url='<img.*?src="(.*?)".*?alt'
	img_urls=re.findall(img_url,request)
	name='<title>(.*?)&nbsp;- 百度文库</title>'
	filename1=re.findall(name,request)
	#打印测试
	print(filename1)
	
	#这里对图片进行遍历，并保存文件
	i=0
	for x in img_urls:
		#print(img_urls)
		url2=x.replace('amp;','').replace("mp;","")
		#print(url2)
		path=r'C:\Users\admin\Desktop\Lanzizi'
		if not os.path.exists(path):
		 	os.mkdir(path)
		os.chdir(path)
		try:
			res=requests.get(url2)
			file=str(i)+'.jpg'
			open(file,"wb").write(res.content)
		except Exception as e:
			continue
		i+=1

	return filename1


#ppt 把下载的图片制成ppt 函数
def make_ppt(filename):  # 如果是ppt， 在上面下载好图片后，这个函数嫁给你图片重新合成ppt
	filenameq=filename[0]+'.pptx'
	pptFile= Presentation() 
	#picFiles=[fn for fn in os.listdir(r"C:\Users\admin\Desktop\baidu") if fn.endswith('.jpg')]
	picFiles=[fn for fn in os.listdir(r"C:\Users\admin\Desktop\Lanzizi") if fn.endswith('.jpg')]
	#print(picFiles)
	#修好后的代码
	for i in range(0,len(picFiles)):
		fn="C:\\Users\\admin\\Desktop\\Lanzizi\\"+str(i)+'.jpg'
		slide= pptFile.slides.add_slide(pptFile.slide_layouts[1])
		slide.shapes.add_picture(fn,Inches(0),Inches(0),Inches(10),Inches(7.5))
		os.remove(fn)	
	#os.rmdir("C:\\Users\\admin\\Desktop\\Lanzizi\\")#删除文件夹
	pptFile.save(filenameq)



#word 文字获取 保存 
def get_bite(result):#word
	path=r'C:\Users\admin\Desktop\Lanzizi'
	if not os.path.exists(path):
		os.mkdir(path)
	os.chdir(path)
	html = etree.HTML(result)
	#print(html)
	strall=''

	res=html.xpath('//*[@class="bd"]//p')
	titlename=html.xpath('/html/head/title')
	filename=titlename[0].text+".doc"
	file=open(filename,"w+",encoding='utf-8')
	#print(titlename[0].text)
	#写入文件出现乱码错误，贼烦
	#
	for i in res:
		#print(i)
		#print(str(i.text))
		file.write(str(i.text))
		strall+=str(i.text)
		#print(type(strall))
		#open("baiduword.doc","wb",encoding='utf-8').write(strall)
	file.close()
		
	
#主函数
def main():
	
	##ppt

	print('''		❤ ❤ ❤ ❤ ❤ ❤ 我从不让我女友打水 ❤ ❤ ❤ ❤ ❤ ❤/
		❤ ❤ ❤ ❤ ❤ 因为兰子打水-------一场空 ❤ ❤ ❤ ❤ ❤/
	 	❤ ❤ ❤ ❤ ❤ ❤ ❤ ❤ ❤ ❤ ❤ ❤-------byYyang ❤ ❤ ❤ ❤''')
	num=int(input("1:ppt，2:word   \n请输入对应数字（回车确定）："))
	if num==1:
		print("百度文库ppt下载器--------byYyang")
		url=input("请输入ppt链接地址（回车确定）：")
		print("稍等，程序正在运行阶段。。。。。")
		#获取网页响应1函数
		request1=get_request(url)
		#下载ppt 图片，文件标题的信息
		filename=down_jpg(request1)

		print("稍等，程序正在运行阶段。。。。。")
		#ppt制作函数
		make_ppt(filename)

		print("下载  ppt    完成。。。。。")
		print("文件保存在  Lanzizi 文件夹里。。。。。")
		print("ppt头和尾有广告信息，记得自行删除修改。")

		time.sleep(5)
	
	elif num==2:
		#word
		#url3='https://wenku.baidu.com/view/4e0ef33c88eb172ded630b1c59eef8c75ebf9571.html?from=search'
		print("百度文库word文字提取器--------byYyang")
		url3=input("请输入word链接地址（回车确定）：")
		print("稍等，程序正在运行阶段。。。。。")
		request2=get_request(url3)
		#print(request2)
		print("稍等，程序正在运行阶段。。。。。")
		get_bite(request2)
		
		print("下载  word    完成。。。。。")
		print("文件保存在  Lanzizi  文件夹里。。。。。")
		print("word  格式有问题，记得自行删除修改。")
		time.sleep(5)

	else:
		print("输入有误！！！！")

	#退出程序
	sys.exit()





	
    
#函数运行位置
if __name__ == '__main__':
	main()


'''

#真实链接  ppt图片
https://wkretype.bdimg.com/retype/zoom/66223f730b1c59eef8c7b425?pn=1&   o=jpg_6&    md5sum=fcec6df7adf998866e304bcbd887aeb6&    sign=0a90881f17&    png=0-36256&    jpg=0-65793
#爬取到的链接
https://wkretype.bdimg.com/retype/zoom/66223f730b1c59eef8c7b425?pn=1&amp;o=jpg_6&amp;md5sum=fcec6df7adf998866e304bcbd887aeb6&amp;sign=0a90881f17&amp;png=0-36256&amp;jpg=0-65793

4  https://wkretype.bdimg.com/retype/zoom/66223f730b1c59eef8c7b425?pn=4&o=jpg_6&md5sum=fcec6df7adf998866e304bcbd887aeb6&sign=0a90881f17&png=51216-62275&jpg=176505-244495
5  https://wkretype.bdimg.com/retype/zoom/66223f730b1c59eef8c7b425?pn=5&o=jpg_6&md5sum=fcec6df7adf998866e304bcbd887aeb6&sign=0a90881f17&png=62276-73254&jpg=244496-283909
6  https://wkretype.bdimg.com/retype/zoom/66223f730b1c59eef8c7b425?pn=6&o=jpg_6&md5sum=fcec6df7adf998866e304bcbd887aeb6&sign=0a90881f17&png=73255-81470&jpg=283910-335515
猜 https://wkretype.bdimg.com/retype/zoom/66223f730b1c59eef8c7b425?pn=7&o=jpg_6&md5sum=fcec6df7adf998866e304bcbd887aeb6&sign=0a90881f17&png=81471-81470&jpg=283910-335515

真https://wkretype.bdimg.com/retype/zoom/66223f730b1c59eef8c7b425?pn=7&o=jpg_6&md5sum=fcec6df7adf998866e304bcbd887aeb6&sign=0a90881f17&png=81471-91541&jpg=335516-372465

#word:


'''
